"""File extraction helpers for PDFs, PPTX files, DOCX files, and plain text."""

from __future__ import annotations

import re
from io import BytesIO
from typing import Iterator

import fitz
from docx import Document
from docx.document import Document as DocxDocument
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract ordered text from a PDF while preserving lightweight structure."""
    document = fitz.open(stream=file_bytes, filetype="pdf")
    pages: list[str] = []

    try:
        for page in document:
            blocks = sorted(
                page.get_text("blocks"),
                key=lambda block: (round(block[1], 1), round(block[0], 1)),
            )
            lines: list[str] = []

            for block in blocks:
                block_text = (block[4] or "").strip()
                if not block_text:
                    continue

                for raw_line in block_text.splitlines():
                    cleaned = _normalize_extracted_line(raw_line)
                    if cleaned:
                        lines.append(cleaned)

            lines = _dedupe_preserve_order(lines)
            if not lines:
                continue

            if _looks_like_heading(lines[0]):
                page_text = f"## {lines[0]}\n" + "\n".join(lines[1:])
            else:
                page_text = "\n".join(lines)

            pages.append(page_text.strip())
    finally:
        document.close()

    return "\n\n".join(section for section in pages if section).strip()


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract slide titles, bullets, and tables from a PowerPoint file."""
    presentation = Presentation(BytesIO(file_bytes))
    slides: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        title_text = ""

        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title_text = _normalize_extracted_line(slide.shapes.title.text)
            if title_text:
                parts.append(f"## Slide {slide_index}: {title_text}")

        for shape in slide.shapes:
            if slide.shapes.title is not None and shape == slide.shapes.title:
                continue

            if getattr(shape, "has_table", False):
                table_text = _table_to_markdown(
                    [[_normalize_extracted_line(cell.text) for cell in row.cells] for row in shape.table.rows]
                )
                if table_text:
                    parts.append(table_text)
                continue

            if not getattr(shape, "has_text_frame", False):
                continue

            for paragraph in shape.text_frame.paragraphs:
                text = _normalize_extracted_line("".join(run.text for run in paragraph.runs) or paragraph.text)
                if not text or text == title_text:
                    continue

                if paragraph.level == 0 and _looks_like_heading(text):
                    parts.append(f"### {text}")
                else:
                    indent = "  " * min(paragraph.level, 3)
                    parts.append(f"{indent}- {text}")

        slide_text = "\n".join(_dedupe_preserve_order(parts)).strip()
        if slide_text:
            slides.append(slide_text)

    return "\n\n".join(slides).strip()


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract paragraphs and tables from a Word document in document order."""
    document = Document(BytesIO(file_bytes))
    blocks: list[str] = []

    for block in _iter_docx_blocks(document):
        if isinstance(block, Paragraph):
            text = _normalize_extracted_line(block.text)
            if not text:
                continue

            style_name = (block.style.name or "").lower() if block.style else ""
            if style_name.startswith("heading"):
                level = _heading_level_from_style(style_name)
                blocks.append(f"{'#' * level} {text}")
            elif "bullet" in style_name or "list" in style_name:
                blocks.append(f"- {text}")
            else:
                blocks.append(text)
            continue

        if isinstance(block, Table):
            rows = [[_normalize_extracted_line(cell.text) for cell in row.cells] for row in block.rows]
            table_text = _table_to_markdown(rows)
            if table_text:
                blocks.append(table_text)

    return "\n\n".join(section for section in blocks if section).strip()


def extract_text_from_txt(file_bytes: bytes) -> str:
    """Decode plain text files with a forgiving fallback."""
    return file_bytes.decode("utf-8", errors="ignore").strip()


def _iter_docx_blocks(document: DocxDocument) -> Iterator[Paragraph | Table]:
    """Yield paragraphs and tables in their original order."""
    parent_elm = document.element.body
    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield Table(child, document)


def _heading_level_from_style(style_name: str) -> int:
    match = re.search(r"heading\s+(\d+)", style_name)
    if not match:
        return 2
    return max(1, min(int(match.group(1)), 4))


def _normalize_extracted_line(line: str) -> str:
    line = line.replace("\xa0", " ")
    line = line.replace("\u2022", "-")
    line = re.sub(r"\s+", " ", line).strip()
    return line


def _looks_like_heading(line: str) -> bool:
    if len(line) > 90:
        return False
    if len(line.split()) > 12:
        return False
    if line.endswith((".", ";", "?", "!")):
        return False
    return bool(re.search(r"[A-Za-z0-9]", line))


def _dedupe_preserve_order(lines: list[str]) -> list[str]:
    deduped: list[str] = []
    previous = ""
    for line in lines:
        if line == previous:
            continue
        deduped.append(line)
        previous = line
    return deduped


def _table_to_markdown(rows: list[list[str]]) -> str:
    cleaned_rows = [[cell.strip() for cell in row] for row in rows if any(cell.strip() for cell in row)]
    if not cleaned_rows:
        return ""

    column_count = max(len(row) for row in cleaned_rows)
    normalized_rows = [row + [""] * (column_count - len(row)) for row in cleaned_rows]
    header = normalized_rows[0]
    divider = ["---"] * column_count
    table_lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(divider) + " |",
    ]
    table_lines.extend("| " + " | ".join(row) + " |" for row in normalized_rows[1:])
    return "\n".join(table_lines)
